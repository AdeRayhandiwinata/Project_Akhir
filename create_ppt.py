from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───
DARK_BG = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x10, 0xB9, 0x81)
ACCENT_LIGHT = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
PURPLE_LIGHT = RGBColor(0xF3, 0xE8, 0xFF)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
TEAL_LIGHT = RGBColor(0xCC, 0xFB, 0xF1)
RED_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)
GREEN_TEXT = RGBColor(0x16, 0xA3, 0x4A)
GRAY_100 = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_600 = RGBColor(0x47, 0x55, 0x69)
GRAY_700 = RGBColor(0x33, 0x40, 0x55)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
HEADER_H = Inches(1.2)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_DARK,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=16, color=TEXT_DARK,
                spacing=8, bullet=""):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet} {item}" if bullet else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return txBox


def add_header(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, DARK_BG)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             title, font_size=32, color=WHITE, bold=True)


def add_stat_card(slide, left, top, width, height, icon, value, label, bg):
    add_shape(slide, left, top, width, height, bg)
    add_text(slide, left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.4),
             icon, font_size=24, color=TEXT_DARK)
    add_text(slide, left + Inches(0.25), top + Inches(0.65), width - Inches(0.5), Inches(0.5),
             value, font_size=24, color=TEXT_DARK, bold=True)
    add_text(slide, left + Inches(0.25), top + Inches(1.15), width - Inches(0.5), Inches(0.35),
             label, font_size=12, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT
circle.line.fill.background()
add_text(slide, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1), "Car",
         font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.0), Inches(11.333), Inches(1),
         "DriveRent", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.9), Inches(11.333), Inches(0.7),
         "Sistem Informasi Rental Mobil", font_size=26, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
         "Pemrograman Web  |  SIB 4A  |  2026", font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: DAFTAR ISI
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.1), Inches(7.5), ACCENT)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
         "Daftar Isi", font_size=34, color=TEXT_DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT)

toc = [
    ("01", "Tentang Proyek"),
    ("02", "Tech Stack"),
    ("03", "Landing Page"),
    ("04", "Armada Mobil"),
    ("05", "Cara Sewa"),
    ("06", "Dashboard Admin"),
    ("07", "Fitur Dashboard"),
    ("08", "Struktur Project"),
    ("09", "Penutup"),
]

for i, (num, title) in enumerate(toc):
    y = Inches(1.6) + Inches(i * 0.6)
    add_text(slide, Inches(1.2), y, Inches(0.7), Inches(0.45),
             num, font_size=20, color=ACCENT, bold=True)
    add_text(slide, Inches(2.1), y + Inches(0.03), Inches(5), Inches(0.4),
             title, font_size=18, color=TEXT_DARK, bold=False)
    if i < len(toc) - 1:
        add_rect(slide, Inches(2.1), y + Inches(0.5), Inches(8), Inches(0.01), GRAY_100)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: TENTANG PROYEK
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Tentang Proyek")

card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(7.2), Inches(5.2), CARD_BG)
add_text(slide, Inches(1.2), Inches(1.9), Inches(6.5), Inches(0.5),
         "DriveRent - Rental Mobil Terpercaya", font_size=22, color=TEXT_DARK, bold=True)
add_rect(slide, Inches(1.2), Inches(2.5), Inches(6), Inches(0.03), ACCENT)

desc = [
    "Aplikasi web rental mobil lengkap dengan landing page dan dashboard admin",
    "Dirancang untuk kemudahan pelanggan dalam menyewa kendaraan",
    "Dashboard admin untuk manajemen armada, transaksi, dan pelanggan",
    "Integrasi WhatsApp untuk pemesanan langsung",
    "Responsive design - dapat diakses dari desktop maupun mobile",
]
add_bullets(slide, Inches(1.2), Inches(2.8), Inches(6.3), Inches(3.5),
            desc, font_size=15, spacing=10, bullet=">")

add_stat_card(slide, Inches(8.5), Inches(1.6), Inches(4), Inches(2.0),
              "Car", "24", "Total Armada", BLUE_LIGHT)
add_stat_card(slide, Inches(8.5), Inches(3.9), Inches(4), Inches(2.0),
              "People", "156", "Pelanggan Aktif", PURPLE_LIGHT)
add_stat_card(slide, Inches(8.5), Inches(6.2), Inches(4), Inches(0.9),
              "Money", "Rp 45JT", "Pendapatan/Bulan", GREEN_LIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: TECH STACK
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Tech Stack")

tech = [
    ("Next.js 16", "React framework dengan App Router, Server Components, dan routing otomatis", BLUE_LIGHT),
    ("React 19", "Library UI modern dengan Server Components dan performa optimal", GREEN_LIGHT),
    ("TypeScript", "Type-safe JavaScript untuk kode yang lebih reliable dan maintainable", BLUE_LIGHT),
    ("Tailwind CSS 4", "Utility-first CSS untuk design responsif dan konsisten", TEAL_LIGHT),
    ("Lucide React", "Library ikon lightweight dengan 1000+ ikon SVG berkualitas", PURPLE_LIGHT),
    ("shadcn/ui", "Komponen UI yang dapat dikustomisasi, built on Radix UI", AMBER_LIGHT),
]

for i, (name, desc, bg) in enumerate(tech):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.6) + Inches(row * 2.7)
    add_shape(slide, x, y, Inches(3.7), Inches(2.3), bg)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5),
             name, font_size=20, color=TEXT_DARK, bold=True)
    add_rect(slide, x + Inches(0.3), y + Inches(0.85), Inches(2.5), Inches(0.03), ACCENT)
    add_text(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(1.0),
             desc, font_size=13, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: LANDING PAGE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Landing Page - DriveRent")

sections = [
    ("Navbar", "Navigasi sticky dengan logo, menu Armada, Cara Sewa, Testimoni, Kontak. Responsive dengan hamburger menu di mobile.", BLUE_LIGHT),
    ("Hero Section", "Headline Perjalanan Nyaman, Harga Terjangkau dengan CTA Lihat Armada dan Chat WhatsApp. Trust signals: Asuransi, 24 Jam.", GREEN_LIGHT),
    ("Car List", "Grid 6 mobil: Avanza, Brio, Innova Reborn, Civic, Fortuner, Xenia. Card dengan harga, spesifikasi, dan tombol Sewa.", PURPLE_LIGHT),
    ("How It Works", "3 langkah mudah: Pilih Mobil > Atur Jadwal > Ambil & Nikmati. Ilustrasi icon dan deskripsi singkat.", AMBER_LIGHT),
    ("Testimoni", "3 review pelanggan dengan quote pengalaman. Menampilkan nama, role, dan cerita pengalaman menyewa.", TEAL_LIGHT),
    ("Footer", "Info kontak: telepon, email, website. Jam operasional dan hak cipta DriveRent.", RED_LIGHT),
]

for i, (name, desc, bg) in enumerate(sections):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.6) + Inches(row * 2.7)
    add_shape(slide, x, y, Inches(3.7), Inches(2.3), bg)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), Inches(3.1), Inches(0.45),
             name, font_size=19, color=TEXT_DARK, bold=True)
    add_rect(slide, x + Inches(0.3), y + Inches(0.75), Inches(2.5), Inches(0.03), ACCENT)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.1), Inches(1.1),
             desc, font_size=12, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: ARMADA MOBIL
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Armada Mobil")

# Table
add_rect(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.55), DARK_BG)
cols = [("Mobil", 1.0, 2.8), ("Tipe", 4.0, 1.5), ("Kursi", 5.7, 1.0),
        ("Bensin", 6.9, 1.2), ("Transmisi", 8.3, 1.5), ("Harga/Hari", 10.0, 2.3)]
for h, x, w in cols:
    add_text(slide, Inches(x), Inches(1.65), Inches(w), Inches(0.45),
             h, font_size=13, color=WHITE, bold=True)

cars = [
    ("Toyota Avanza", "MPV", "7", "Bensin", "Manual", "Rp 250.000"),
    ("Honda Brio", "Hatchback", "5", "Bensin", "Automatic", "Rp 200.000"),
    ("Toyota Innova Reborn", "MPV Premium", "7", "Diesel", "Automatic", "Rp 400.000"),
    ("Honda Civic", "Sedan", "5", "Bensin", "Automatic", "Rp 350.000"),
    ("Toyota Fortuner", "SUV", "7", "Diesel", "Automatic", "Rp 500.000"),
    ("Daihatsu Xenia", "MPV", "7", "Bensin", "Manual", "Rp 230.000"),
]

for i, car in enumerate(cars):
    y = Inches(2.25) + Inches(i * 0.55)
    bg = CARD_BG if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.5), bg)
    for j, (_, x, w) in enumerate(cols):
        add_text(slide, Inches(x), y + Inches(0.07), Inches(w), Inches(0.36),
                 car[j], font_size=12, color=TEXT_DARK, bold=(j == 0))

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
         "Populer: Toyota Avanza, Innova Reborn, Toyota Fortuner", font_size=13, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(6.25), Inches(11), Inches(0.4),
         "Semua kendaraan dirawat secara berkala dan dilengkapi asuransi untuk perjalanan yang aman.",
         font_size=13, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: CARA SEWA
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Cara Sewa - 3 Langkah Mudah")

steps = [
    ("01", "Pilih Mobil", "Telusuri armada kami dan pilih kendaraan yang paling sesuai dengan kebutuhan perjalanan Anda.", BLUE_LIGHT),
    ("02", "Atur Jadwal", "Tentukan tanggal dan durasi sewa. Hubungi kami via WhatsApp untuk konfirmasi ketersediaan.", GREEN_LIGHT),
    ("03", "Ambil & Nikmati", "Ambil kendaraan di lokasi yang telah disepakati atau kami antar ke lokasi Anda.", PURPLE_LIGHT),
]

for i, (num, title, desc, bg) in enumerate(steps):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.8)
    add_shape(slide, x, y, Inches(3.7), Inches(4.8), bg)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.35), y + Inches(0.4), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_text(slide, x + Inches(1.35), y + Inches(0.6), Inches(1), Inches(0.6),
             num, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.3), y + Inches(1.7), Inches(3.1), Inches(0.5),
             title, font_size=22, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(1.2), y + Inches(2.25), Inches(1.3), Inches(0.03), ACCENT)
    add_text(slide, x + Inches(0.4), y + Inches(2.5), Inches(2.9), Inches(1.8),
             desc, font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

for i in range(2):
    x = Inches(4.65) + Inches(i * 4.1)
    add_text(slide, x, Inches(3.8), Inches(0.4), Inches(0.5),
             ">", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: DASHBOARD ADMIN
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Dashboard Admin - Balikpapan Trans")

# Sidebar
add_shape(slide, Inches(0.8), Inches(1.6), Inches(3.2), Inches(5.4), DARK_BG)
add_text(slide, Inches(1.1), Inches(1.8), Inches(2.7), Inches(0.5),
         "Balikpapan Trans", font_size=18, color=WHITE, bold=True)
add_rect(slide, Inches(1.1), Inches(2.35), Inches(2.7), Inches(0.02), GRAY_600)

menu = ["Dashboard", "Daftar Armada", "Transaksi", "Pelanggan"]
for i, label in enumerate(menu):
    y = Inches(2.6) + Inches(i * 0.55)
    bg = ACCENT if i == 0 else GRAY_700
    add_shape(slide, Inches(1.1), y, Inches(2.7), Inches(0.45), bg)
    add_text(slide, Inches(1.3), y + Inches(0.05), Inches(2.4), Inches(0.35),
             label, font_size=13, color=WHITE, bold=(i == 0))

# Content area
add_text(slide, Inches(4.5), Inches(1.7), Inches(8), Inches(0.5),
         "Dashboard Overview", font_size=20, color=TEXT_DARK, bold=True)

stats = [
    ("Car", "24", "Total Armada", "+2 bulan ini", BLUE_LIGHT),
    ("List", "8", "Transaksi Hari Ini", "+3 dari kemarin", GREEN_LIGHT),
    ("People", "156", "Pelanggan Aktif", "+12 bulan ini", PURPLE_LIGHT),
    ("Money", "Rp 45JT", "Pendapatan", "+18% dari bulan lalu", AMBER_LIGHT),
]

for i, (icon, value, label, change, bg) in enumerate(stats):
    x = Inches(4.5) + Inches(i * 2.15)
    y = Inches(2.3)
    add_shape(slide, x, y, Inches(2.0), Inches(1.6), bg)
    add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(1.7), Inches(0.35),
             icon, font_size=18, color=TEXT_DARK)
    add_text(slide, x + Inches(0.15), y + Inches(0.45), Inches(1.7), Inches(0.35),
             value, font_size=18, color=TEXT_DARK, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.85), Inches(1.7), Inches(0.25),
             label, font_size=10, color=TEXT_MUTED)
    add_text(slide, x + Inches(0.15), y + Inches(1.15), Inches(1.7), Inches(0.25),
             change, font_size=9, color=GREEN_TEXT)

# Table
add_text(slide, Inches(4.5), Inches(4.2), Inches(8), Inches(0.4),
         "Transaksi Terakhir", font_size=15, color=TEXT_DARK, bold=True)

add_rect(slide, Inches(4.5), Inches(4.7), Inches(8.3), Inches(0.4), GRAY_100)
th = [("Pelanggan", 4.7), ("Mobil", 6.7), ("Tanggal", 8.7), ("Status", 10.5)]
for h, x in th:
    add_text(slide, Inches(x), Inches(4.73), Inches(1.8), Inches(0.35),
             h, font_size=11, color=TEXT_MUTED, bold=True)

rows = [
    ("Budi Santoso", "Toyota Avanza", "05 Jul 2026", "Selesai"),
    ("Siti Rahayu", "Honda Brio", "05 Jul 2026", "Berlangsung"),
    ("Ahmad Fauzi", "Innova Reborn", "04 Jul 2026", "Selesai"),
]
for i, (name, car, date, status) in enumerate(rows):
    y = Inches(5.2) + Inches(i * 0.45)
    bg = CARD_BG if i % 2 == 0 else WHITE
    add_rect(slide, Inches(4.5), y, Inches(8.3), Inches(0.4), bg)
    vals = [name, car, date, status]
    for val, (_, x) in zip(vals, th):
        add_text(slide, Inches(x), y + Inches(0.05), Inches(1.8), Inches(0.3),
                 val, font_size=11, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: FITUR DASHBOARD
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Fitur Dashboard Admin")

features = [
    ("Manajemen Armada", [
        "Daftar semua kendaraan (24 unit)",
        "Status: Tersedia / Disewa",
        "Filter berdasarkan tipe & status",
        "CRUD: Tambah, Edit, Hapus armada",
        "Info: nama, tipe, tahun, plat, warna, harga",
    ], BLUE_LIGHT),
    ("Manajemen Transaksi", [
        "Riwayat seluruh transaksi rental",
        "Status: Selesai, Berlangsung, Menunggu",
        "Filter berdasarkan status & tanggal",
        "Detail: pelanggan, mobil, durasi, total",
        "Ringkasan total pendapatan",
    ], GREEN_LIGHT),
    ("Manajemen Pelanggan", [
        "Database pelanggan lengkap (156 akun)",
        "Info: nama, email, telepon, alamat",
        "Total riwayat sewa per pelanggan",
        "Filter berdasarkan status aktif/tidak",
        "Pelanggan baru: +12 bulan ini",
    ], PURPLE_LIGHT),
]

for i, (title, items, bg) in enumerate(features):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.6)
    add_shape(slide, x, y, Inches(3.7), Inches(5.3), bg)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5),
             title, font_size=20, color=TEXT_DARK, bold=True)
    add_rect(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.1), Inches(0.03), ACCENT)
    add_bullets(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(3.8),
                items, font_size=14, spacing=10, bullet=">")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: STRUKTUR PROJECT
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Struktur Project")

# File tree
add_shape(slide, Inches(0.8), Inches(1.6), Inches(6.5), Inches(5.4), DARK_BG)
tree = [
    "rental-mobil/",
    "  app/",
    "    page.tsx",
    "    layout.tsx",
    "    globals.css",
    "    components/",
    "      Sidebar.tsx",
    "    dashboard/",
    "      page.tsx",
    "      layout.tsx",
    "      armada/page.tsx",
    "      transaksi/page.tsx",
    "      pelanggan/page.tsx",
    "  components/",
    "    Navbar.tsx",
    "    Hero.tsx",
    "    CarList.tsx",
    "    HowItWorks.tsx",
    "    Testimonials.tsx",
    "    Footer.tsx",
    "    ui/ (shadcn)",
    "  package.json",
    "  tsconfig.json",
    "  next.config.ts",
]
for i, line in enumerate(tree):
    y = Inches(1.8) + Inches(i * 0.2)
    is_dir = line.strip().endswith("/") or line.strip().startswith("app") or line.strip().startswith("components") or line.strip().startswith("dashboard")
    color = ACCENT_LIGHT if is_dir else SLATE_400
    indent = len(line) - len(line.lstrip())
    add_text(slide, Inches(1.1) + Inches(indent * 0.08), y, Inches(5.5), Inches(0.22),
             line, font_size=10, color=color, font_name="Consolas")

# Explanation
add_text(slide, Inches(7.8), Inches(1.7), Inches(5), Inches(0.5),
         "Penjelasan Struktur", font_size=20, color=TEXT_DARK, bold=True)
add_rect(slide, Inches(7.8), Inches(2.25), Inches(2), Inches(0.03), ACCENT)

explanations = [
    "app/ - App Router (Next.js 16)",
    "  page.tsx - Landing page utama",
    "  dashboard/ - Admin panel",
    "  components/ - Komponen reusable",
    "",
    "components/ - Landing page UI",
    "  Navbar, Hero, CarList, dll",
    "  ui/ - Komponen shadcn/ui",
    "",
    "Config: package.json, tsconfig, next.config",
]
add_bullets(slide, Inches(7.8), Inches(2.5), Inches(5), Inches(4),
            explanations, font_size=14, spacing=6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: KEUNGGULAN
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Keunggulan Aplikasi")

advantages = [
    ("Modern & Cepat", "Built dengan Next.js 16 & React 19 untuk performa optimal dan loading time cepat.", BLUE_LIGHT),
    ("Responsive Design", "Tampilan optimal di desktop, tablet, dan mobile. Akses dari perangkat apapun.", GREEN_LIGHT),
    ("Type-Safe", "TypeScript memastikan kode bebas error runtime dan lebih mudah di-maintain.", PURPLE_LIGHT),
    ("Easy Booking", "Integrasi WhatsApp untuk pemesanan instan. Pelanggan langsung chat tanpa ribet.", AMBER_LIGHT),
    ("Admin Dashboard", "Panel lengkap untuk kelola armada, transaksi, dan pelanggan dalam satu tempat.", TEAL_LIGHT),
    ("Modern UI/UX", "Desain profesional dengan Tailwind CSS & shadcn/ui. Clean dan intuitif.", RED_LIGHT),
]

for i, (title, desc, bg) in enumerate(advantages):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.6) + Inches(row * 2.7)
    add_shape(slide, x, y, Inches(3.7), Inches(2.3), bg)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), Inches(3.1), Inches(0.45),
             title, font_size=19, color=TEXT_DARK, bold=True)
    add_rect(slide, x + Inches(0.3), y + Inches(0.75), Inches(2.5), Inches(0.03), ACCENT)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.1), Inches(1.0),
             desc, font_size=13, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: PENUTUP
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

add_text(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(1),
         "Terima Kasih", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.7),
         "DriveRent - Rental Mobil Terpercaya", font_size=22, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.667), Inches(4.2), Inches(2), Inches(2))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT
circle.line.fill.background()
add_text(slide, Inches(5.667), Inches(4.7), Inches(2), Inches(1),
         "Q&A", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.6), Inches(11.333), Inches(0.4),
         "Pemrograman Web  |  SIB 4A  |  2026", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT)


# ─── Save ───
output_path = os.path.join(os.path.dirname(__file__), "DriveRent_Presentasi_v2.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
